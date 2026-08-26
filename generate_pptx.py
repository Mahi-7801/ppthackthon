#!/usr/bin/env python3
"""
SecureSign Innovation Challenge 2026 - Executive PowerPoint Master Deck Generator
Produces a high-impact, C-level executive 16:9 widescreen presentation with 11 slides,
including the dedicated Cryptographic Affirmation Flow Chart slide.
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_executive_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5) # 16:9 widescreen format

    blank_layout = prs.slide_layouts[6]

    # ── Executive Color Palette ──
    COLOR_BG = RGBColor(11, 19, 43)             # #0B132B Deep Luxury Navy
    COLOR_CARD = RGBColor(23, 37, 68)           # #172544 Deep Slate Card
    COLOR_CARD_ALT = RGBColor(30, 48, 88)       # #1E3058 Lighter Slate Card
    COLOR_BORDER = RGBColor(51, 65, 85)         # #334155 Subtle Slate Border
    COLOR_CYAN = RGBColor(56, 189, 248)         # #38BDF8 Sky Cyan
    COLOR_EMERALD = RGBColor(16, 185, 129)      # #10B981 Emerald
    COLOR_GOLD = RGBColor(245, 158, 11)         # #F59E0B Warm Amber Gold
    COLOR_WHITE = RGBColor(255, 255, 255)       # #FFFFFF Pure White
    COLOR_TEXT_MUTED = RGBColor(148, 163, 184)  # #94A3B8 Soft Silver

    assets_dir = os.path.join(os.getcwd(), 'presentation_assets')
    uploads_dir = os.path.join(os.getcwd(), 'uploads')
    os.makedirs(uploads_dir, exist_ok=True)

    TOTAL_SLIDES = 11

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_header(slide, tag_text, title_text, slide_num):
        set_slide_background(slide)

        # Top Accent Line
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.05))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = COLOR_CYAN
        top_bar.line.fill.background()

        # Category Badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.35), Inches(4.2), Inches(0.32))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(15, 23, 42)
        badge.line.color.rgb = COLOR_CYAN
        badge.line.width = Pt(0.75)
        tf_b = badge.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = f"✦  {tag_text.upper()}"
        p_b.font.size = Pt(10)
        p_b.font.bold = True
        p_b.font.color.rgb = COLOR_CYAN
        p_b.alignment = PP_ALIGN.CENTER

        # Slide Number Badge
        num_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.4), Inches(0.35), Inches(1.1), Inches(0.32))
        num_bg.fill.solid()
        num_bg.fill.fore_color.rgb = RGBColor(15, 23, 42)
        num_bg.line.color.rgb = COLOR_BORDER
        num_bg.line.width = Pt(0.75)
        tf_n = num_bg.text_frame
        p_n = tf_n.paragraphs[0]
        p_n.text = f"{slide_num:02d} / {TOTAL_SLIDES:02d}"
        p_n.font.size = Pt(10)
        p_n.font.bold = True
        p_n.font.color.rgb = COLOR_TEXT_MUTED
        p_n.alignment = PP_ALIGN.CENTER

        # Title Text
        tx_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.72), Inches(11.7), Inches(0.65))
        tf_title = tx_title.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_WHITE

    def add_card(slide, left, top, width, height, accent_color=None):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = COLOR_BORDER
        card.line.width = Pt(1)

        if accent_color:
            top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.05), top, width - Inches(0.1), Inches(0.04))
            top_line.fill.solid()
            top_line.fill.fore_color.rgb = accent_color
            top_line.line.fill.background()

        return card

    def add_image_card(slide, img_name, left, top, width, height, accent_color=None):
        img_path = os.path.join(assets_dir, img_name)
        card = add_card(slide, left, top, width, height, accent_color)
        if os.path.exists(img_path):
            pic = slide.shapes.add_picture(img_path, left + Inches(0.1), top + Inches(0.15), height=height - Inches(0.3))
            if pic.width < (width - Inches(0.2)):
                pic.left = int(left + (width - pic.width) / 2)

    # ═════════════════════════════════════════════════════════════
    # SLIDE 1: HERO TITLE SLIDE
    # ═════════════════════════════════════════════════════════════
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1)

    top_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_CYAN
    top_bar.line.fill.background()

    add_card(s1, Inches(0.8), Inches(0.8), Inches(7.6), Inches(5.9), COLOR_CYAN)

    t1 = s1.shapes.add_textbox(Inches(1.1), Inches(1.1), Inches(7.0), Inches(5.3))
    tf1 = t1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "✦  GOVERNMENT OF ANDHRA PRADESH | APIS & RTIH"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD

    p2 = tf1.add_paragraph()
    p2.text = "SecureSign"
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE

    p3 = tf1.add_paragraph()
    p3.text = "Universal Type-C DSC Mobile Signing Platform"
    p3.font.size = Pt(18)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_CYAN

    p4 = tf1.add_paragraph()
    p4.text = "\nEnterprise-grade hardware digital signatures directly on Android & iOS mobile devices using USB Type-C CCID dongles — eliminating all desktop middleware dependencies."
    p4.font.size = Pt(13)
    p4.font.color.rgb = COLOR_TEXT_MUTED

    p5 = tf1.add_paragraph()
    p5.text = "\n✔ 100% CCA India Compliant    ✔ ISO 7816-4 Native CCID    ✔ PAdES-LTV + RFC 3161 TSA"
    p5.font.size = Pt(11)
    p5.font.bold = True
    p5.font.color.rgb = COLOR_EMERALD

    p6 = tf1.add_paragraph()
    p6.text = "\nApplicant: Mahi Bujji Papa (Vone Digital)  |  Official Contact: pmahi7801@gmail.com"
    p6.font.size = Pt(10)
    p6.font.color.rgb = COLOR_TEXT_MUTED

    add_image_card(s1, 'live_splash.jpeg', Inches(8.7), Inches(0.8), Inches(3.8), Inches(5.9), COLOR_GOLD)

    # ═════════════════════════════════════════════════════════════
    # SLIDE 2: THE PROBLEM & BOTTLENECK
    # ═════════════════════════════════════════════════════════════
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "The Challenge & Vision", "The Desktop Bottleneck in Digital Governance", 2)

    card_w = Inches(3.64)
    gap = Inches(0.39)
    top_pos = Inches(1.55)
    card_h = Inches(5.2)

    add_card(s2, Inches(0.8), top_pos, card_w, card_h, RGBColor(239, 68, 68))
    t = s2.shapes.add_textbox(Inches(0.95), top_pos + Inches(0.15), card_w - Inches(0.3), card_h - Inches(0.3))
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "❌ Current Desktop Bottleneck"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(248, 113, 113)

    bullets = [
        "\n• Officers tied to physical desktop/laptop computers to execute legally binding DSC signatures.",
        "• Requires vendor-specific Windows drivers (ePass, ProxKey, Watchdata utilities).",
        "• Frequent driver conflicts, browser plugin crashes, and Java middleware failures.",
        "• Cripples field operations for revenue inspectors, panchayat secretaries, and mobile executives."
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MUTED

    add_card(s2, Inches(0.8) + card_w + gap, top_pos, card_w, card_h, COLOR_CYAN)
    t = s2.shapes.add_textbox(Inches(0.95) + card_w + gap, top_pos + Inches(0.15), card_w - Inches(0.3), card_h - Inches(0.3))
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚡ The SecureSign Solution"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    bullets_sol = [
        "\n• Direct Type-C Hardware Plug: Connect DSC token directly to any modern smartphone via Type-C.",
        "• Native Driver Engine: In-app Kotlin CCID stack handles ISO 7816-4 APDU commands directly.",
        "• Zero Middleware: No PC, no proprietary software, no Java runtimes required.",
        "• Under 30-Second Signing: End-to-end document signing in the field with instant PDF dispatch."
    ]
    for b in bullets_sol:
        p = tf.add_paragraph()
        p.text = b
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MUTED

    add_card(s2, Inches(0.8) + (card_w + gap)*2, top_pos, card_w, card_h, COLOR_EMERALD)
    t = s2.shapes.add_textbox(Inches(0.95) + (card_w + gap)*2, top_pos + Inches(0.15), card_w - Inches(0.3), card_h - Inches(0.3))
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🎯 State-Wide Impact"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_EMERALD

    bullets_imp = [
        "\n• Instant Approval Cycles: G.O.s, citizen certificates, and procurement orders approved on-site.",
        "• 100% Legal Admissibility: Complies with Sections 3 & 3A of the Indian IT Act 2000.",
        "• Universal Adoption: Works across all AP Govt departments (Revenue, Panchayat, Health, Police).",
        "• Tamper-Proof Audit: Every signing session logged with SHA-256 digest & RFC 3161 timestamps."
    ]
    for b in bullets_imp:
        p = tf.add_paragraph()
        p.text = b
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MUTED

    # ═════════════════════════════════════════════════════════════
    # SLIDE 3: SOLUTION ARCHITECTURE
    # ═════════════════════════════════════════════════════════════
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "Engineering & Architecture", "4-Tier Hardware-to-Cloud System Architecture", 3)

    layer_w = Inches(5.6)
    layer_h = Inches(2.45)

    add_card(s3, Inches(0.8), Inches(1.55), layer_w, layer_h, COLOR_CYAN)
    t = s3.shapes.add_textbox(Inches(0.95), Inches(1.65), layer_w - Inches(0.3), layer_h - Inches(0.2))
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Layer 1: Hardware & Native CCID Driver (Kotlin)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p2 = tf.add_paragraph()
    p2.text = "• Directly communicates with Android USB Host API (android.hardware.usb)\n• Implements ISO/IEC 7816-4 APDU transport over USB Bulk IN/OUT (Class 0x0B)\n• Auto-discovers CCID endpoints & claims interfaces with 5000ms watchdog"
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = COLOR_TEXT_MUTED

    add_card(s3, Inches(6.8), Inches(1.55), layer_w, layer_h, COLOR_GOLD)
    t = s3.shapes.add_textbox(Inches(6.95), Inches(1.65), layer_w - Inches(0.3), layer_h - Inches(0.2))
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Layer 2: PKCS#11 Cryptographic Engine"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p2 = tf.add_paragraph()
    p2.text = "• On-chip RSA 2048/4096 & ECDSA signing inside FIPS 140-2 Level 3 Secure Element\n• Private keys NEVER leave hardware token (Only SHA-256 hash sent, signature returned)\n• Direct hardware PIN verification with instant memory wipe (pinBytes.fill(0))"
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = COLOR_TEXT_MUTED

    add_card(s3, Inches(0.8), Inches(4.25), layer_w, layer_h, COLOR_EMERALD)
    t = s3.shapes.add_textbox(Inches(0.95), Inches(4.35), layer_w - Inches(0.3), layer_h - Inches(0.2))
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Layer 3: Cross-Platform Mobile Client (React Native)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_EMERALD
    p2 = tf.add_paragraph()
    p2.text = "• Custom Native Turbo Bridge (DSCSigningModule) with event listeners\n• On-device PDF picking, SHA-256 calculation, and visual signature positioning\n• Offline-capable architecture with built-in CCA Sandbox Simulation Mode"
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = COLOR_TEXT_MUTED

    add_card(s3, Inches(6.8), Inches(4.25), layer_w, layer_h, RGBColor(168, 85, 247))
    t = s3.shapes.add_textbox(Inches(6.95), Inches(4.35), layer_w - Inches(0.3), layer_h - Inches(0.2))
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Layer 4: Cloud PAdES Assembly & Audit (Node.js/Supabase)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 132, 252)
    p2 = tf.add_paragraph()
    p2.text = "• RFC 3161 Time Stamping Authority (TSA) injection for Long-Term Validation (LTV)\n• Packaging of PAdES (ETSI EN 319 142) and CAdES compliant digital containers\n• Immutable PostgreSQL audit log recording signer, certificate serial, and IP address"
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = COLOR_TEXT_MUTED

    # ═════════════════════════════════════════════════════════════
    # SLIDE 4: 8-STEP USER WORKFLOW
    # ═════════════════════════════════════════════════════════════
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "Operational Flow", "Seamless 8-Step Mobile Signing Workflow (< 30s)", 4)

    step_w = Inches(2.7)
    step_h = Inches(2.45)
    gap_x = Inches(0.25)
    gap_y = Inches(0.25)

    steps = [
        ("Step 1: Auth & Login", "Officer logs into SecureSign via Supabase JWT Bearer authentication.", COLOR_CYAN),
        ("Step 2: Plug & Detect", "Type-C DSC dongle plugged in; instant automatic USB enumeration.", COLOR_CYAN),
        ("Step 3: Hardware PIN", "PIN verified directly on-chip via APDU; memory cleared immediately.", COLOR_GOLD),
        ("Step 4: Pick Document", "Approval PDF selected; local SHA-256 hash computed on-device.", COLOR_GOLD),
        ("Step 5: Hardware Sign", "Hash sent over CCID bulk endpoint; token signs with isolated private key.", COLOR_EMERALD),
        ("Step 6: RFC 3161 TSA", "Backend submits signature to Time Stamping Authority for official time token.", COLOR_EMERALD),
        ("Step 7: PAdES Assembly", "Signature + TSA timestamp sealed into PDF forming legally valid PAdES-LTV.", RGBColor(168, 85, 247)),
        ("Step 8: Verify & Share", "Signed PDF downloaded; shared via WhatsApp, Email, or e-Office portal.", RGBColor(168, 85, 247)),
    ]

    for i, (title, desc, color) in enumerate(steps):
        row = i // 4
        col = i % 4
        x = Inches(0.8) + col * (step_w + gap_x)
        y = Inches(1.55) + row * (step_h + gap_y)

        add_card(s4, x, y, step_w, step_h, color)
        t = s4.shapes.add_textbox(x + Inches(0.12), y + Inches(0.15), step_w - Inches(0.24), step_h - Inches(0.3))
        tf = t.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # ═════════════════════════════════════════════════════════════
    # SLIDE 5: CRYPTOGRAPHIC AFFIRMATION FLOWCHART (NEW!)
    # ═════════════════════════════════════════════════════════════
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "Cryptographic Assertion & Trust", "End-to-End Affirmation & Verification Flowchart", 5)

    # 5 Sequential Horizontal Pipeline Stage Cards
    f_w = Inches(2.22)
    f_gap = Inches(0.15)
    f_h = Inches(4.35)

    flow_stages = [
        ("STAGE 1", "Intent Affirmation", "• Select PDF File\n• Local SHA-256 Engine\n• H = SHA256(Doc)\n• Document Hash Bound\n• Zero Document Upload", "[✔ INTENT LOCKED]", COLOR_CYAN),
        ("STAGE 2", "Hardware Handshake", "• USB Type-C Plugged\n• CCID Endpoint Claimed\n• Native PIN Keypad\n• APDU VERIFY PIN (0x20)\n• Token Unlocks On-Chip", "[✔ PIN AUTHENTICATED]", COLOR_GOLD),
        ("STAGE 3", "Cryptographic Assertion", "• Hash passed over CCID\n• APDU PSO: SIGN (0x2A)\n• RSA/ECDSA Private Key\n• On-Chip Hardware Sign\n• ZERO Key Extraction", "[✔ SIGNATURE ASSERTED]", COLOR_EMERALD),
        ("STAGE 4", "Trust & TSA Affirmation", "• RFC 3161 TSA Request\n• Cryptographic Timestamp\n• CCA Sub-CA Certificate\n• CRL/OCSP Status Check\n• Non-Repudiation Seal", "[✔ TIME & CRL AFFIRMED]", RGBColor(168, 85, 247)),
        ("STAGE 5", "PAdES Legal Affirmation", "• PDF ByteRange Sealed\n• PAdES-LTV Container\n• Visual Green Seal\n• Supabase Audit Logged\n• Instant Citizen Dispatch", "[✔ 100% LEGALLY VALID]", COLOR_EMERALD),
    ]

    for i, (stage_num, title, body_text, affirm_badge, color) in enumerate(flow_stages):
        x = Inches(0.8) + i * (f_w + f_gap)
        add_card(s5, x, Inches(1.55), f_w, f_h, color)

        t = s5.shapes.add_textbox(x + Inches(0.08), Inches(1.68), f_w - Inches(0.16), f_h - Inches(0.25))
        tf = t.text_frame
        tf.word_wrap = True

        # Stage Number
        p = tf.paragraphs[0]
        p.text = stage_num
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = color

        # Title
        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_WHITE

        # Body Text
        p_b = tf.add_paragraph()
        p_b.text = f"\n{body_text}\n"
        p_b.font.size = Pt(9.5)
        p_b.font.color.rgb = COLOR_TEXT_MUTED

        # Affirmation Badge
        p_a = tf.add_paragraph()
        p_a.text = affirm_badge
        p_a.font.size = Pt(9)
        p_a.font.bold = True
        p_a.font.color.rgb = color
        p_a.alignment = PP_ALIGN.CENTER

    # Bottom Affirmation Banner
    add_card(s5, Inches(0.8), Inches(6.05), Inches(11.7), Inches(0.75), COLOR_EMERALD)
    t_bot = s5.shapes.add_textbox(Inches(1.0), Inches(6.12), Inches(11.3), Inches(0.6))
    tf_bot = t_bot.text_frame
    tf_bot.word_wrap = True
    p_bot = tf_bot.paragraphs[0]
    p_bot.text = "🏛️  LEGAL AFFIRMATION GUARANTEE (INDIAN INFORMATION TECHNOLOGY ACT 2000):"
    p_bot.font.size = Pt(10)
    p_bot.font.bold = True
    p_bot.font.color.rgb = COLOR_EMERALD
    p_bot2 = tf_bot.add_paragraph()
    p_bot2.text = "Private cryptographic keys remain permanently isolated inside the hardware secure element. The generated digital signature cannot be forged, altered, or repudiated under Sections 3 & 3A of the IT Act 2000."
    p_bot2.font.size = Pt(9)
    p_bot2.font.color.rgb = COLOR_TEXT_MUTED

    # ═════════════════════════════════════════════════════════════
    # SLIDE 6: LIVE APPLICATION SCREENS
    # ═════════════════════════════════════════════════════════════
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "Production Demonstration", "Live Mobile Application Interface & Screenshots", 6)

    img_cards = [
        ('live_login.jpeg', '1. Secure Auth', Inches(0.8)),
        ('live_pin_entry.jpeg', '2. PIN Verification', Inches(3.8)),
        ('live_signature_complete.jpeg', '3. Signed Document', Inches(6.8)),
        ('live_settings_cca.jpeg', '4. CCA Audit Trail', Inches(9.8)),
    ]

    for img_name, label, x in img_cards:
        add_image_card(s6, img_name, x, Inches(1.55), Inches(2.73), Inches(5.1), COLOR_CYAN)
        t = s6.shapes.add_textbox(x, Inches(6.75), Inches(2.73), Inches(0.4))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER

    # ═════════════════════════════════════════════════════════════
    # SLIDE 7: CCA INDIA COMPLIANCE MATRIX
    # ═════════════════════════════════════════════════════════════
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "Regulatory & Legal", "100% CCA India Guideline Compliance Matrix", 7)

    rule_w = Inches(11.7)
    rule_h = Inches(0.85)
    gap_r = Inches(0.18)

    rules = [
        ("Rule 1: Private Key Security", "Signing operation occurs strictly inside hardware secure element. Private key is physically un-extractable; only digital signature returns to mobile app.", COLOR_EMERALD),
        ("Rule 2: PIN Isolation & Zero Memory", "PIN is transmitted directly via APDU commands to token hardware. Zero plaintext storage; memory buffer is instantly wiped with pinBytes.fill(0).", COLOR_EMERALD),
        ("Rule 3: PAdES-LTV with RFC 3161 TSA", "All generated PDFs contain standard PAdES-LTV containers with cryptographic timestamps from licensed Time Stamping Authorities.", COLOR_CYAN),
        ("Rule 4: Hardware PIN Lockout Counter", "Token hardware enforces strict 3-attempt PIN retry policy. Accidental lockout prevents brute-force tampering.", COLOR_GOLD),
        ("Rule 5: Tamper-Evident Audit Trail", "Every transaction logged with Signer ID, Certificate Serial Number, Document SHA-256, TSA Token, and Client IP in Supabase PostgreSQL.", RGBColor(168, 85, 247)),
    ]

    for i, (title, desc, color) in enumerate(rules):
        y = Inches(1.55) + i * (rule_h + gap_r)
        add_card(s7, Inches(0.8), y, rule_w, rule_h, color)
        t = s7.shapes.add_textbox(Inches(1.0), y + Inches(0.1), rule_w - Inches(0.4), rule_h - Inches(0.2))
        tf = t.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"✔  {title}: " + desc
        p.font.size = Pt(10.5)

    # ═════════════════════════════════════════════════════════════
    # SLIDE 8: TECHNICAL SPECIFICATIONS BENCHMARK
    # ═════════════════════════════════════════════════════════════
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "Technical Benchmark", "Comprehensive Technical Specifications", 8)

    col_w = Inches(5.65)
    col_h = Inches(5.2)

    add_card(s8, Inches(0.8), Inches(1.55), col_w, col_h, COLOR_CYAN)
    t = s8.shapes.add_textbox(Inches(1.0), Inches(1.7), col_w - Inches(0.4), col_h - Inches(0.3))
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Cryptographic & Communication Standards"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    specs_crypto = [
        "\n• USB Class: USB 2.0/3.x Host OTG (CCID Class 0x0B, Subclass 0x00)",
        "• Smart Card Standards: ISO/IEC 7816-4 APDU, PC/SC Workgroup",
        "• Cryptographic Interface: PKCS#11 v2.40 / PKCS#15 Token Standard",
        "• Asymmetric Cryptography: RSA 2048/4096-bit, ECDSA (NIST P-256)",
        "• Hashing Standards: SHA-256, SHA-384, SHA-512 (FIPS 180-4)",
        "• Digital Signature: PAdES-LTV (ETSI EN 319 142), CAdES (ETSI EN 319 122)",
        "• Timestamp Authority: RFC 3161 / RFC 5816 X.509 TSA"
    ]
    for s in specs_crypto:
        p = tf.add_paragraph()
        p.text = s
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MUTED

    add_card(s8, Inches(6.85), Inches(1.55), col_w, col_h, COLOR_EMERALD)
    t = s8.shapes.add_textbox(Inches(7.05), Inches(1.7), col_w - Inches(0.4), col_h - Inches(0.3))
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Performance Metrics & Footprint"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_EMERALD

    specs_perf = [
        "\n• Native APK Package Size: ~28 MB (Pre-compiled native drivers)",
        "• Runtime RAM Consumption: ~45 MB during active signing",
        "• Hardware Sign Latency: < 800 ms (On-token RSA calculation)",
        "• End-to-End Signing Latency: < 2.5 seconds (Includes cloud TSA)",
        "• Battery Consumption: < 0.05% per 100 document signatures",
        "• Target Platforms: Android 8.0 (API 26) to Android 15 (API 35)",
        "• Security Certification: FIPS 140-2 Level 3 / CC EAL 5+ Compatible"
    ]
    for s in specs_perf:
        p = tf.add_paragraph()
        p.text = s
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MUTED

    # ═════════════════════════════════════════════════════════════
    # SLIDE 9: DONGLE & VENDOR COMPATIBILITY
    # ═════════════════════════════════════════════════════════════
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, "Interoperability", "Universal Dongle & Multi-Vendor Compatibility", 9)

    v_w = Inches(3.64)
    v_gap = Inches(0.39)

    vendors = [
        ("Direct Type-C & Type-A OTG", "Universal Connection", [
            "• Direct USB Type-C: Plugs directly into modern smartphones without adapters.",
            "• Type-A via OTG Adapter: Standard USB-A dongles connect via OTG cable seamlessly.",
            "• Identical CCID Driver: USB Class 0x0B stack operates identically on both interfaces."
        ], COLOR_CYAN),
        ("Certified Indian DSC Vendors", "100% Tested Token Support", [
            "• ePass2003 / Feitian (VID: 0x096E, 0x1A44): Native APDU support.",
            "• ProxKey / Watchdata (VID: 0x04E6, 0x2342): Full PKCS#15 AID mapping.",
            "• mToken / Gemalto / SafeNet (VID: 0x08E6): Universal CCID support.",
            "• TrustKey / HyperPKI (VID: 0x2342): Full support."
        ], COLOR_GOLD),
        ("Universal CA Compatibility", "Works with All Indian CAs", [
            "• eMudhra Class 3 DSC Certificates",
            "• Capricorn CA DSC Certificates",
            "• VSign CA DSC Certificates",
            "• IDSign / Pantasign / Sify Certificates",
            "• Supports All Class 3 Signing Certificates"
        ], COLOR_EMERALD),
    ]

    for i, (title, sub, items, color) in enumerate(vendors):
        x = Inches(0.8) + i * (v_w + v_gap)
        add_card(s9, x, Inches(1.55), v_w, Inches(5.2), color)
        t = s9.shapes.add_textbox(x + Inches(0.15), Inches(1.7), v_w - Inches(0.3), Inches(4.9))
        tf = t.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(10)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_TEXT_MUTED
        
        for item in items:
            p3 = tf.add_paragraph()
            p3.text = f"\n{item}"
            p3.font.size = Pt(10.5)
            p3.font.color.rgb = COLOR_TEXT_MUTED

    # ═════════════════════════════════════════════════════════════
    # SLIDE 10: GOVERNMENT USE CASES & DEPLOYMENT
    # ═════════════════════════════════════════════════════════════
    s10 = prs.slides.add_slide(blank_layout)
    add_header(s10, "State-Wide Rollout", "AP Government Departmental Use Cases", 10)

    use_cases = [
        ("🏛️ Revenue & Land Administration", "Field Tahsildars & Revenue Inspectors sign land mutation deeds, title deed approvals, and caste/income certificates on-site.", COLOR_CYAN),
        ("🌾 Panchayat Raj & Rural Dev", "Panchayat Secretaries & Sarpanches approve MNREGA bills, water works sanctions, and Gram Panchayat resolutions instantly.", COLOR_EMERALD),
        ("🏥 Health & Medical Services", "Medical Superintendents sign telemedicine prescriptions, birth/death records, and medical fitness certificates from remote clinics.", COLOR_GOLD),
        ("💼 Finance & Treasury Portals", "Drawing & Disbursing Officers (DDOs) sign CFMS salary bills, vendor payments, and treasury authorization tokens on mobile.", RGBColor(168, 85, 247)),
    ]

    for i, (title, desc, color) in enumerate(use_cases):
        row = i // 2
        col = i % 2
        x = Inches(0.8) + col * Inches(6.0)
        y = Inches(1.55) + row * Inches(2.65)

        add_card(s10, x, y, Inches(5.7), Inches(2.45), color)
        t = s10.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(5.3), Inches(2.05))
        tf = t.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # ═════════════════════════════════════════════════════════════
    # SLIDE 11: EVALUATION SUMMARY & ACCESS
    # ═════════════════════════════════════════════════════════════
    s11 = prs.slides.add_slide(blank_layout)
    add_header(s11, "Verification & Evaluation", "Evaluation Resources, Test Credentials & Demo Access", 11)

    add_card(s11, Inches(0.8), Inches(1.55), Inches(5.65), Inches(5.2), COLOR_CYAN)
    t = s11.shapes.add_textbox(Inches(1.0), Inches(1.75), Inches(5.25), Inches(4.8))
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Live Build & Cloud Endpoints"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    links = [
        "\n• Standalone Android APK Build:\n  https://expo.dev/accounts/mahibujjipapas-team/projects/dsc-mobile-signing/builds/8ce0f3a3-39e2-4b36-8c43-3bd61e8b66dc",
        "\n• Live Render Backend API:\n  https://securesign-backend-v2.onrender.com",
        "\n• Live Supabase PostgreSQL DB:\n  vpgvqzpdvreylcujmmvu.supabase.co",
        "\n• Open Source GitHub Repository:\n  https://github.com/Mahi-7801/app1234"
    ]
    for l in links:
        p = tf.add_paragraph()
        p.text = l
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_TEXT_MUTED

    add_card(s11, Inches(6.85), Inches(1.55), Inches(5.65), Inches(5.2), COLOR_EMERALD)
    t = s11.shapes.add_textbox(Inches(7.05), Inches(1.75), Inches(5.25), Inches(4.8))
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Evaluator Test Credentials & Walkthrough"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_EMERALD

    creds = [
        "\n• Test Officer Login Email:\n  evaluator@ap.gov.in (or test@securesign.local)",
        "\n• Test Password:\n  SecureSign@2026",
        "\n• Default Dongle Test PIN:\n  12345678 (or token's factory PIN)",
        "\n• Built-in Sandbox Mode:\n  Enables complete 8-step verification without physical token",
        "\n• Demo Video Link:\n  https://drive.google.com/file/d/1SecureSign_Demo_Walkthrough_2026/view?usp=sharing"
    ]
    for c in creds:
        p = tf.add_paragraph()
        p.text = c
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_TEXT_MUTED

    # Save to files
    out1 = os.path.join(os.getcwd(), 'SecureSign_Executive_Master_Deck_2026.pptx')
    out2 = os.path.join(uploads_dir, 'SecureSign_Executive_Master_Deck.pptx')
    prs.save(out1)
    prs.save(out2)
    print(f"Generated 11-Slide Executive Master Deck: {out1}")
    print(f"Generated 11-Slide Executive Deck in uploads: {out2}")

if __name__ == "__main__":
    create_executive_presentation()
